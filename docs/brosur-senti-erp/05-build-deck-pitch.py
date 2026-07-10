#!/usr/bin/env python3
"""
Deck Pitch PPT Senti ERP — 16:9, ~12 slide, untuk meeting presentasi client.
Output: 05-deck-pitch-senti-erp.pptx

Dibangun dengan python-pptx. Edit copy di DATA di bawah, jalankan ulang.
  python3 05-build-deck-pitch.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import copy

# ---------- palette ----------
BLUE   = RGBColor(0x25, 0x63, 0xEB)
BLUEH  = RGBColor(0x1D, 0x4E, 0xD8)
BLUEDK = RGBColor(0x0B, 0x1A, 0x3D)
BLUEBG = RGBColor(0xEE, 0xF4, 0xFF)
TX     = RGBColor(0x0F, 0x17, 0x2A)
MUT    = RGBColor(0x5B, 0x64, 0x77)
BD     = RGBColor(0xE6, 0xE8, 0xEC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x15, 0x80, 0x3D)
REDBG  = RGBColor(0xFE, 0xF2, 0xF2)
REDTX  = RGBColor(0x7F, 0x1D, 0x1D)
GRNBG  = RGBColor(0xEC, 0xFD, 0xF5)
GRNTX  = RGBColor(0x06, 0x5F, 0x46)
AMBER  = RGBColor(0xCA, 0x8A, 0x04)

FONT = "Inter"

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9 widescreen
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]

# ---------- helpers ----------
def slide(bg=WHITE):
    s = prs.slides.add_slide(blank)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.line.fill.background()
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.shadow.inherit = False
    return s

def txt(s, x, y, w, h, runs, size=18, color=TX, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.15, space_after=0):
    """runs: str OR list of (text, dict-overrides)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if isinstance(runs, str):
        runs = [(runs, {})]
    first = True
    for text, ov in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = ov.get("align", align)
        p.line_spacing = ov.get("line_spacing", line_spacing)
        p.space_after = Pt(ov.get("space_after", space_after))
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = ov.get("font", font)
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.italic = ov.get("italic", False)
        c = ov.get("color", color)
        f.color.rgb = c
        first = False
    return tb

def rect(s, x, y, w, h, fill=BLUE, line=None, radius=None):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if radius is not None:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp

def card(s, x, y, w, h, fill=WHITE, border=BD, radius=0.06):
    return rect(s, x, y, w, h, fill=fill, line=border, radius=radius)

def page_chrome(s, idx, total, dark=False):
    """footer with logo + page number."""
    col = WHITE if dark else TX
    mcol = RGBColor(0x93,0xC5,0xFD) if dark else MUT
    # logo mark
    mk = rect(s, Inches(0.55), Inches(7.02), Inches(0.32), Inches(0.32),
              fill=BLUE if not dark else WHITE, radius=0.25)
    txt(s, Inches(0.92), Inches(7.0), Inches(2), Inches(0.36),
        "Senti ERP", size=12, bold=True, color=col, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(11.0), Inches(7.0), Inches(1.8), Inches(0.36),
        f"{idx:02d} / {total:02d}", size=11, color=mcol,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    # top accent line
    rect(s, 0, 0, SW, Inches(0.08), fill=BLUE)

def eyebrow(s, text, x=Inches(0.7), y=Inches(0.7), color=BLUE):
    txt(s, x, y, Inches(8), Inches(0.3), text.upper(),
        size=12, bold=True, color=color)

TOTAL = 12

# ============================================================
# SLIDE 1 — COVER
# ============================================================
s = slide(BLUEDK)
rect(s, 0, 0, SW, SH, fill=BLUEDK)
# gradient-ish overlay band
rect(s, 0, Inches(2.6), SW, Inches(0.06), fill=BLUE)
mk = rect(s, Inches(0.7), Inches(0.6), Inches(0.5), Inches(0.5), fill=WHITE, radius=0.2)
txt(s, Inches(1.3), Inches(0.62), Inches(4), Inches(0.5), "Senti ERP",
    size=20, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
eyebrow(s, "Pitch Deck · 2026", y=Inches(1.7), color=RGBColor(0x93,0xC5,0xFD))
txt(s, Inches(0.7), Inches(2.85), Inches(12), Inches(2.4),
    [("ERP Lengkap untuk Bisnis Indonesia.",
      {"size":50, "bold":True, "color":WHITE, "line_spacing":1.05})])
txt(s, Inches(0.7), Inches(4.7), Inches(11), Inches(1.0),
    "Akuntansi, stok, penjualan, dan produksi dalam satu sistem berbasis web. "
    "Harga flat, pengguna tak terbatas, implementasi 2 minggu.",
    size=20, color=RGBColor(0xC7,0xD2,0xFE), line_spacing=1.35)
# badges
bx = Inches(0.7)
for b in ["2.500+ perusahaan", "Data di Indonesia", "99,9% uptime"]:
    w = Inches(2.3)
    pill = rect(s, bx, Inches(5.9), w, Inches(0.42), fill=RGBColor(0x16,0x29,0x4D), radius=0.5)
    txt(s, bx, Inches(5.9), w, Inches(0.42), b, size=12, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bx = Emu(int(bx) + int(w) + Inches(0.2))
txt(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
    "PT Senti Teknologi Indonesia  ·  021-5051-5105",
    size=13, color=RGBColor(0x93,0xC5,0xFD))

# ============================================================
# SLIDE 2 — AGENDA
# ============================================================
s = slide(); page_chrome(s, 2, TOTAL)
eyebrow(s, "Agenda")
txt(s, Inches(0.7), Inches(1.05), Inches(12), Inches(0.8),
    "Apa yang akan kita bahas", size=34, bold=True)
items = [
    ("01", "Masalah", "Pain operasional yang lazim dijumpai."),
    ("02", "Solusi", "Apa itu Senti ERP dan kenapa beda."),
    ("03", "Modul & Industri", "Fitur dan segmen yang dilayani."),
    ("04", "Harga", "Model harga flat, tanpa per user."),
    ("05", "Implementasi", "Bagaimana onboarding berjalan."),
    ("06", "Demo & Q&A", "Jadwalkan demo & tanya jawab."),
]
x0, y0, gap = Inches(0.7), Inches(2.4), Inches(0.25)
cw = Inches((13.333 - 1.4 - 0.5)/3)  # 3 cols
ch = Inches(1.25)
for i,(num,t,d) in enumerate(items):
    col = i % 3; row = i // 3
    x = Emu(int(x0) + col*(int(cw)+int(gap)))
    y = Emu(int(y0) + row*(int(ch)+int(gap)))
    card(s, x, y, cw, ch, fill=BLUEBG)
    txt(s, x+Inches(0.25), y+Inches(0.18), Inches(0.8), Inches(0.5),
        num, size=22, bold=True, color=BLUEH)
    txt(s, x+Inches(1.05), y+Inches(0.2), cw-Inches(1.2), Inches(0.4),
        t, size=17, bold=True, color=TX)
    txt(s, x+Inches(1.05), y+Inches(0.62), cw-Inches(1.2), Inches(0.55),
        d, size=12, color=MUT, line_spacing=1.2)

# ============================================================
# SLIDE 3 — PROBLEM
# ============================================================
s = slide(); page_chrome(s, 3, TOTAL)
eyebrow(s, "Masalah")
txt(s, Inches(0.7), Inches(1.05), Inches(12), Inches(0.8),
    "Operasional yang menyita waktu & uang", size=32, bold=True)
txt(s, Inches(0.7), Inches(1.85), Inches(11.5), Inches(0.5),
    "Empat pain yang kami dengar berulang dari pemilik usaha:",
    size=16, color=MUT)
pains = [
    ("📉", "Tutup buku makan 10 hari", "Tim keuangan kejar setoran; angka sering berubah sampai akhir bulan."),
    ("📦", "Stok tak cocok dengan pembukuan", "Fisik gudang beda dengan catatan Excel — selisih sulit ditelusuri."),
    ("🔁", "Entri ganda antar tim", "Penjualan input di satu tempat, keuangan input lagi di tempat lain."),
    ("💸", "Software mahal & per user", "Tambah karyawan = tambah biaya; jadi tidak scale dengan bisnis."),
]
y = Inches(2.6)
for ico, t, d in pains:
    card(s, Inches(0.7), y, Inches(11.93), Inches(0.95), fill=WHITE, border=BD)
    rect(s, Inches(0.7), y, Inches(0.08), Inches(0.95), fill=BLUE)
    txt(s, Inches(0.95), y+Inches(0.18), Inches(0.7), Inches(0.6),
        ico, size=26, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.75), y+Inches(0.14), Inches(4.5), Inches(0.4),
        t, size=17, bold=True, color=TX)
    txt(s, Inches(1.75), y+Inches(0.5), Inches(9.8), Inches(0.4),
        d, size=13, color=MUT)
    y = Emu(int(y) + Inches(1.05))

# ============================================================
# SLIDE 4 — SOLUTION (Sebelum / Sesudah)
# ============================================================
s = slide(); page_chrome(s, 4, TOTAL)
eyebrow(s, "Solusi")
txt(s, Inches(0.7), Inches(1.05), Inches(12), Inches(0.8),
    "Bisnis Anda, setelah Senti ERP", size=32, bold=True)
# two columns
colw = Inches(5.85); colh = Inches(4.2); y0 = Inches(2.2)
# before
card(s, Inches(0.7), y0, colw, colh, fill=REDBG, border=RGBColor(0xFC,0xCA,0xCA))
txt(s, Inches(0.95), y0+Inches(0.2), colw, Inches(0.4), "SEBELUM",
    size=14, bold=True, color=REDTX)
bp = ["Tutup buku 10 hari","Stok tidak cocok","Entri ganda antar tim",
      "Bayar per user","Server sendiri rawan mati"]
yy = Emu(int(y0) + Inches(0.75))
for it in bp:
    txt(s, Inches(0.95), yy, colw, Inches(0.5),
        [("✗  ",{"color":REDTX,"bold":True}),(it,{"color":REDTX})],
        size=16, line_spacing=1.2)
    yy = Emu(int(yy)+Inches(0.6))
# after
x2 = Inches(6.78)
card(s, x2, y0, colw, colh, fill=GRNBG, border=RGBColor(0xA7,0xF3,0xD0))
txt(s, x2+Inches(0.25), y0+Inches(0.2), colw, Inches(0.4), "SESUDAH",
    size=14, bold=True, color=GRNTX)
ap = ["Tutup buku 2 hari","Real-time, auto kartu stok","Satu data, semua modul",
      "Harga flat, user tak terbatas","Cloud + backup harian"]
yy = Emu(int(y0) + Inches(0.75))
for it in ap:
    txt(s, x2+Inches(0.25), yy, colw, Inches(0.5),
        [("✓  ",{"color":GREEN,"bold":True}),(it,{"color":GRNTX})],
        size=16, line_spacing=1.2)
    yy = Emu(int(yy)+Inches(0.6))
txt(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.4),
    "Satu sistem menyatukan data → angka lebih cepat, lebih akurat, lebih murah.",
    size=14, bold=True, color=BLUEH, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5 — KENAPA (5 USP)
# ============================================================
s = slide(); page_chrome(s, 5, TOTAL)
eyebrow(s, "Kenapa Senti ERP")
txt(s, Inches(0.7), Inches(1.05), Inches(12), Inches(0.8),
    "Lima keunggulan yang dirasakan langsung", size=30, bold=True)
usps = [
    ("🌐","Berbasis Web","Akses dari browser mana pun, tanpa server sendiri."),
    ("∞","User Tak Terbatas","Satu harga untuk seluruh tim, berapa pun jumlahnya."),
    ("⚡","Real-Time","Setiap transaksi langsung tercermin di laporan."),
    ("🔒","Aman & Teraudit","Enkripsi, backup harian, jejak audit lengkap."),
    ("🤝","Didampingi Tim Lokal","Setup sampai berjalan — implementasi 2 minggu."),
]
cw = Inches(2.36); ch = Inches(3.4); gap = Inches(0.2)
x0 = Inches(0.7); y0 = Inches(2.3)
for i,(ico,t,d) in enumerate(usps):
    x = Emu(int(x0) + i*(int(cw)+int(gap)))
    card(s, x, y0, cw, ch, fill=WHITE, border=BD)
    rect(s, x, y0, cw, Inches(0.08), fill=BLUE)
    txt(s, x+Inches(0.2), y0+Inches(0.35), cw-Inches(0.4), Inches(0.9),
        ico, size=40, align=PP_ALIGN.CENTER)
    txt(s, x+Inches(0.2), y0+Inches(1.5), cw-Inches(0.4), Inches(0.8),
        t, size=18, bold=True, color=TX, align=PP_ALIGN.CENTER, line_spacing=1.1)
    txt(s, x+Inches(0.2), y0+Inches(2.35), cw-Inches(0.4), Inches(1.0),
        d, size=12.5, color=MUT, align=PP_ALIGN.CENTER, line_spacing=1.3)

# ============================================================
# SLIDE 6 — MODUL
# ============================================================
s = slide(); page_chrome(s, 6, TOTAL)
eyebrow(s, "Satu data, semua modul")
txt(s, Inches(0.7), Inches(1.05), Inches(12), Inches(0.8),
    "8 modul terintegrasi, tanpa entri ganda", size=30, bold=True)
mods = [
    ("MD","Master Data","Pelanggan, vendor, item, akun."),
    ("FIN","Akuntansi & GL","Jurnal otomatis, neraca, laba rugi."),
    ("INV","Inventory","Kartu stok multi-gudang, opname."),
    ("PUR","Pembelian","PO, penerimaan, hutang."),
    ("SLS","Penjualan","Faktur, piutang, harga bertingkat."),
    ("MFG","Produksi","BOM, perintah kerja, HPP."),
    ("FA","Aset Tetap","Penyusutan, pelacakan aset."),
    ("POS","POS & Retail","Kasir nyambung ke stok & akuntansi."),
]
cw = Inches(2.95); ch = Inches(1.7); gx = Inches(0.2); gy = Inches(0.2)
x0 = Inches(0.7); y0 = Inches(2.25)
for i,(code,nm,d) in enumerate(mods):
    col = i % 4; row = i // 4
    x = Emu(int(x0) + col*(int(cw)+int(gx)))
    y = Emu(int(y0) + row*(int(ch)+int(gy)))
    card(s, x, y, cw, ch, fill=BLUEBG)
    txt(s, x+Inches(0.25), y+Inches(0.18), Inches(1.0), Inches(0.5),
        code, size=20, bold=True, color=BLUEH)
    txt(s, x+Inches(0.25), y+Inches(0.6), cw-Inches(0.4), Inches(0.4),
        nm, size=15, bold=True, color=TX)
    txt(s, x+Inches(0.25), y+Inches(1.0), cw-Inches(0.4), Inches(0.6),
        d, size=11.5, color=MUT, line_spacing=1.2)
txt(s, Inches(0.7), Inches(6.35), Inches(12), Inches(0.4),
    "Semua modul berbagi satu database — input sekali, terpakai di semua sisi.",
    size=13, color=MUT, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 7 — PER INDUSTRI
# ============================================================
s = slide(); page_chrome(s, 7, TOTAL)
eyebrow(s, "Solusi per industri")
txt(s, Inches(0.7), Inches(1.05), Inches(12), Inches(0.8),
    "Konfigurasi siap untuk model bisnis Anda", size=30, bold=True)
inds = [
    ("🏭","Manufaktur","BOM, perintah kerja, dan biaya produksi (HPP).","MFG + FIN"),
    ("📦","Distribusi","Multi-gudang, harga bertingkat, armada pengiriman.","INV + PUR"),
    ("☕","Retail & F&B","POS terhubung langsung ke stok dan akuntansi.","POS + INV + SLS"),
    ("📋","Jasa & Proyek","Biaya per proyek dan penagihan bertahap.","FIN + SLS"),
]
y = Inches(2.3)
for ico,t,d,mm in inds:
    card(s, Inches(0.7), y, Inches(11.93), Inches(0.95), fill=WHITE, border=BD)
    txt(s, Inches(0.95), y+Inches(0.18), Inches(0.7), Inches(0.6),
        ico, size=24, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.75), y+Inches(0.14), Inches(3.2), Inches(0.4),
        t, size=17, bold=True, color=TX)
    txt(s, Inches(1.75), y+Inches(0.5), Inches(7.5), Inches(0.4),
        d, size=13, color=MUT)
    pill = rect(s, Inches(10.6), y+Inches(0.3), Inches(1.85), Inches(0.36),
                fill=BLUEBG, radius=0.5)
    txt(s, Inches(10.6), y+Inches(0.3), Inches(1.85), Inches(0.36),
        mm, size=11, bold=True, color=BLUEH, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y = Emu(int(y) + Inches(1.05))

# ============================================================
# SLIDE 8 — SOCIAL PROOF
# ============================================================
s = slide(BLUEDK); page_chrome(s, 8, TOTAL, dark=True)
eyebrow(s, "Bukan janji — sudah berjalan", color=RGBColor(0x93,0xC5,0xFD))
txt(s, Inches(0.7), Inches(1.05), Inches(12), Inches(0.8),
    "Angka yang bisa diperiksa", size=34, bold=True, color=WHITE)
stats = [("2.500+","perusahaan pengguna"),
         ("15+","industri terlayani"),
         ("99,9%","uptime 12 bulan"),
         ("4,8/5","rating kepuasan")]
cw = Inches(2.85); gap = Inches(0.2); x0 = Inches(0.7); y0 = Inches(2.6)
for i,(b,l) in enumerate(stats):
    x = Emu(int(x0) + i*(int(cw)+int(gap)))
    card(s, x, y0, cw, Inches(2.4), fill=RGBColor(0x16,0x29,0x4D),
         border=RGBColor(0x2A,0x3B,0x66))
    txt(s, x+Inches(0.2), y0+Inches(0.5), cw-Inches(0.4), Inches(1.0),
        b, size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, x+Inches(0.2), y0+Inches(1.55), cw-Inches(0.4), Inches(0.6),
        l, size=14, color=RGBColor(0xC7,0xD2,0xFE), align=PP_ALIGN.CENTER, line_spacing=1.2)
# testimonial
card(s, Inches(0.7), Inches(5.3), Inches(11.93), Inches(1.3),
     fill=RGBColor(0x16,0x29,0x4D), border=RGBColor(0x2A,0x3B,0x66))
rect(s, Inches(0.7), Inches(5.3), Inches(0.08), Inches(1.3), fill=BLUE)
txt(s, Inches(1.0), Inches(5.45), Inches(11.4), Inches(1.0),
    [("\u201CTutup buku bulanan dari 10 hari jadi 2 hari. Angkanya sekarang bisa dipercaya.\u201D",
      {"size":17,"color":WHITE,"italic":True,"line_spacing":1.3})])
txt(s, Inches(1.0), Inches(6.2), Inches(11.4), Inches(0.35),
    "Budi Santoso — Direktur Operasional, PT Prima Karya Logem · Surabaya",
    size=12, color=RGBColor(0x93,0xC5,0xFD))

# ============================================================
# SLIDE 9 — HARGA
# ============================================================
s = slide(); page_chrome(s, 9, TOTAL)
eyebrow(s, "Harga flat & sederhana")
txt(s, Inches(0.7), Inches(1.05), Inches(12), Inches(0.8),
    "Tanpa biaya per user", size=32, bold=True)
txt(s, Inches(0.7), Inches(1.85), Inches(12), Inches(0.4),
    "Semua paket termasuk pengguna tak terbatas. Tambah tim tanpa tambah biaya.",
    size=15, color=MUT)
plans = [
    ("Starter", "Rp 750rb", "/bln", False,
     ["Modul FIN, INV, SLS, PUR","Pengguna tak terbatas","1 entitas usaha","Dukungan email & chat"]),
    ("Business", "Rp 1,5jt", "/bln", True,
     ["Semua modul, + MFG & POS","Pengguna tak terbatas","3 entitas usaha","Multi-gudang & cabang","Dukungan prioritas"]),
    ("Enterprise", "Kustom", "", False,
     ["Semua fitur Business","Entitas tak terbatas","API & integrasi khusus","SLA & manajer akun"]),
]
cw = Inches(3.83); gap = Inches(0.22); x0 = Inches(0.7); y0 = Inches(2.5)
ch = Inches(3.7)
for i,(nm,pr,unit,pop,feats) in enumerate(plans):
    x = Emu(int(x0) + i*(int(cw)+int(gap)))
    fill = BLUEBG if pop else WHITE
    border = BLUE if pop else BD
    card(s, x, y0, cw, ch, fill=fill, border=border)
    if pop:
        rect(s, x, y0, cw, Inches(0.08), fill=BLUE)
        pill = rect(s, Emu(int(x)+int(cw)-Inches(1.5)), Emu(int(y0)-Inches(0.18)),
                    Inches(1.4), Inches(0.36), fill=BLUE, radius=0.5)
        txt(s, Emu(int(x)+int(cw)-Inches(1.5)), Emu(int(y0)-Inches(0.18)),
            Inches(1.4), Inches(0.36), "PALING POPULER",
            size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x+Inches(0.3), y0+Inches(0.25), cw-Inches(0.6), Inches(0.4),
        nm, size=18, bold=True, color=TX)
    txt(s, x+Inches(0.3), y0+Inches(0.75), cw-Inches(0.6), Inches(0.7),
        [(pr,{"size":32,"bold":True,"color":BLUEH}),
         (unit,{"size":14,"color":MUT})], line_spacing=1.0)
    yy = Emu(int(y0) + Inches(1.65))
    for f in feats:
        txt(s, x+Inches(0.3), yy, cw-Inches(0.6), Inches(0.35),
            [("✓  ",{"color":GREEN,"bold":True}),(f,{"color":TX,"size":12.5})],
            line_spacing=1.2)
        yy = Emu(int(yy)+Inches(0.38))
txt(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.4),
    "Semua paket: coba gratis 30 hari · tanpa kartu kredit · implementasi < 2 minggu",
    size=13, bold=True, color=BLUEH, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 10 — IMPLEMENTASI
# ============================================================
s = slide(); page_chrome(s, 10, TOTAL)
eyebrow(s, "Implementasi")
txt(s, Inches(0.7), Inches(1.05), Inches(12), Inches(0.8),
    "Dari nol sampai jalan dalam 2 minggu", size=30, bold=True)
steps = [
    ("Minggu 1","Setup & Migrasi","Master data & saldo awal diimpor dari template Excel. Tim kami verifikasi bersama Anda."),
    ("Minggu 1-2","Pelatihan Tim","Sesi training per peran: keuangan, gudang, sales, kasir. Materi + rekaman."),
    ("Minggu 2","Go-Live Bertahap","Mulai dari satu modul/unit, lalu roll-out. Dampingan langsung saat transaksi pertama."),
    ("Berkelanjutan","Dukungan","Email, chat, dan (paket Business+) dukungan prioritas + pendampingan berkelanjutan."),
]
y = Inches(2.3)
for badge,t,d in steps:
    card(s, Inches(0.7), y, Inches(11.93), Inches(0.95), fill=WHITE, border=BD)
    pill = rect(s, Inches(0.95), y+Inches(0.28), Inches(1.8), Inches(0.38),
                fill=BLUEBG, radius=0.5)
    txt(s, Inches(0.95), y+Inches(0.28), Inches(1.8), Inches(0.38),
        badge, size=11, bold=True, color=BLUEH, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(3.0), y+Inches(0.14), Inches(4), Inches(0.4),
        t, size=17, bold=True, color=TX)
    txt(s, Inches(3.0), y+Inches(0.5), Inches(9.2), Inches(0.4),
        d, size=13, color=MUT)
    y = Emu(int(y) + Inches(1.05))
txt(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.4),
    "Migrasi dari spreadsheet? Kami periksa hasil impor Excel bersama Anda — bukan self-service.",
    size=13, color=MUT, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 11 — DIFFERENSIATOR vs KOMPETITOR
# ============================================================
s = slide(); page_chrome(s, 11, TOTAL)
eyebrow(s, "Kenapa bukan yang lain")
txt(s, Inches(0.7), Inches(1.05), Inches(12), Inches(0.8),
    "Empat pembeda yang kami pertahankan", size=30, bold=True)
diffs = [
    ("Harga flat, bukan per user","Kompetitor umumnya nagih per user → makin lama makin mahal. Kami tidak."),
    ("Implementasi cepat (1–2 minggu)","Dengan dampingan tim lokal — bukan self-service yang menyatroni."),
    ("Data tersimpan di Indonesia","Kepatuhan UU PDP, latency, kedaulatan data."),
    ("Satu data, semua modul","Tanpa entri ganda, tanpa silo antar departemen."),
]
y = Inches(2.35)
for i,(t,d) in enumerate(diffs):
    card(s, Inches(0.7), y, Inches(11.93), Inches(0.85), fill=WHITE, border=BD)
    rect(s, Inches(0.7), y, Inches(0.08), Inches(0.85), fill=BLUE)
    txt(s, Inches(1.0), y+Inches(0.12), Inches(0.6), Inches(0.6),
        f"{i+1:02d}", size=22, bold=True, color=BLUEH)
    txt(s, Inches(1.75), y+Inches(0.12), Inches(5), Inches(0.4),
        t, size=16, bold=True, color=TX)
    txt(s, Inches(1.75), y+Inches(0.46), Inches(10), Inches(0.35),
        d, size=12.5, color=MUT)
    y = Emu(int(y) + Inches(0.95))

# ============================================================
# SLIDE 12 — CTA / NEXT STEP
# ============================================================
s = slide(BLUEDK); page_chrome(s, 12, TOTAL, dark=True)
eyebrow(s, "Langkah berikutnya", color=RGBColor(0x93,0xC5,0xFD))
txt(s, Inches(0.7), Inches(1.4), Inches(12), Inches(1.5),
    [("Jadwalkan Demo Gratis.",{"size":46,"bold":True,"color":WHITE,"line_spacing":1.1})])
txt(s, Inches(0.7), Inches(3.0), Inches(11), Inches(1.0),
    "20 menit demo berbasis kasus bisnis Anda. Lihat langsung modul yang relevan, "
    "tanya jawab, lalu coba gratis 30 hari.",
    size=18, color=RGBColor(0xC7,0xD2,0xFE), line_spacing=1.4)
# CTA box
card(s, Inches(0.7), Inches(4.4), Inches(11.93), Inches(1.5),
     fill=RGBColor(0x16,0x29,0x4D), border=RGBColor(0x2A,0x3B,0x66))
txt(s, Inches(1.0), Inches(4.55), Inches(5.5), Inches(0.5),
    "📞 Telepon Sales", size=14, bold=True, color=RGBColor(0x93,0xC5,0xFD))
txt(s, Inches(1.0), Inches(4.95), Inches(5.5), Inches(0.7),
    "021-5051-5105", size=28, bold=True, color=WHITE)
txt(s, Inches(6.8), Inches(4.55), Inches(5.5), Inches(0.5),
    "🌐 PT Senti Teknologi Indonesia", size=14, bold=True, color=RGBColor(0x93,0xC5,0xFD))
txt(s, Inches(6.8), Inches(4.95), Inches(5.5), Inches(0.7),
    "Coba Gratis 30 Hari", size=28, bold=True, color=WHITE)
txt(s, Inches(0.7), Inches(6.25), Inches(12), Inches(0.4),
    "Terima kasih.  ·  Implementasi < 2 minggu · Tanpa biaya per pengguna · Data di Indonesia.",
    size=13, color=RGBColor(0x93,0xC5,0xFD), align=PP_ALIGN.CENTER)

out = "/opt/sentient-factory/docs/brosur-senti-erp/05-deck-pitch-senti-erp.pptx"
prs.save(out)
print("OK:", out)
print("Slides:", len(prs.slides._sldIdLst))
